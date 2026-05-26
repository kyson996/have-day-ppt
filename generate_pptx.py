#!/usr/bin/env python3
"""Generate HAVE DAY PPT as .pptx file"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
INK = RGBColor(0x0A, 0x0A, 0x0B)
INK_LIGHT = RGBColor(0x6B, 0x72, 0x80)
INK_TERTIARY = RGBColor(0x9C, 0xA3, 0xAF)
PAPER = RGBColor(0xF1, 0xEF, 0xEA)
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
ACCENT_INDIGO = RGBColor(0x4F, 0x46, 0xE5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_dark_slide():
    """Dark background slide"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = INK
    return slide

def add_light_slide():
    """Light background slide"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PAPER
    return slide

def tb(slide, left, top, width, height, text, font_size=Pt(18),
       color=None, bold=False, align=PP_ALIGN.LEFT, font_name='PingFang SC',
       italic=False):
    """Add a text box"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.font.color.rgb = color or INK
    p.alignment = align
    return tf

def add_rich_tb(slide, left, top, width, height, parts, align=PP_ALIGN.LEFT):
    """Add text box with mixed formatting. parts = [(text, {props})]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, props in parts:
        run = p.add_run()
        run.text = text
        if 'size' in props: run.font.size = props['size']
        if 'bold' in props: run.font.bold = props['bold']
        if 'color' in props: run.font.color.rgb = props['color']
        if 'name' in props: run.font.name = props['name']
        if 'italic' in props: run.font.italic = props['italic']
    return tf

def add_line(slide, left, top, width, color=None, opacity=0.3):
    """Add a thin line"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color or INK
    line.line.fill.background()
    return line

def add_kicker(slide, text, is_dark=False):
    tb(slide, Inches(0.8), Inches(0.6), Inches(10), Inches(0.4),
       text.upper(), Pt(9), INK_TERTIARY if not is_dark else RGBColor(0x99,0x99,0x99),
       font_name='SF Mono')

def add_footer(slide, left_text, right_text, is_dark=False):
    color = RGBColor(0x99,0x99,0x99) if is_dark else INK_TERTIARY
    tb(slide, Inches(0.8), Inches(6.7), Inches(6), Inches(0.3),
       left_text, Pt(8), color, font_name='SF Mono')
    tb(slide, Inches(8), Inches(6.7), Inches(5), Inches(0.3),
       right_text, Pt(8), color, font_name='SF Mono', align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 01 · Cover (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "HAVE DAY · Product Launch", True)
add_footer(slide, "AI 驱动的电子日记", "— 2026 —", True)
tb(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.3),
   "AI 驱动的电子日记", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(1.2),
   "HAVE DAY", Pt(96), RGBColor(0xF1,0xEF,0xEA), bold=True)
tb(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.6),
   "记录每一天，美好每一天", Pt(28), RGBColor(0xBB,0xBB,0xBB))
tb(slide, Inches(0.8), Inches(4.6), Inches(8), Inches(0.8),
   "在这个 AI 时代，用最自然的方式记录生活——说一句话、拍一张照，AI 帮你写成有温度的日记。",
   Pt(16), RGBColor(0xAA,0xAA,0xAA))
tb(slide, Inches(0.8), Inches(5.6), Inches(8), Inches(0.4),
   "产品发布会 · 2026  ·  HAVE DAY Team", Pt(11), RGBColor(0x88,0x88,0x88),
   font_name='SF Mono')

# ============================================================
# SLIDE 02 · Data Numbers (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "Act I · 为什么需要 HAVE DAY", False)
add_footer(slide, "数据来源 · 综合研究", "02 / 15", False)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "一个你可能不想承认的事实", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.7),
   "我们每天都在遗忘", Pt(48), INK, bold=True)
tb(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(0.4),
   "人类大脑不是硬盘，记忆会褪色、扭曲、消失。", Pt(16), INK_LIGHT)

stats = [
    ("Forgetting Curve", "70%", "一天内忘记70%的新信息"),
    ("Daily Moments", "60K+", "每人每天约6万个念头"),
    ("Photo Taken", "5.3亿", "全球每天拍摄的照片"),
    ("Journaling Rate", "<8%", "坚持写日记的人不到8%"),
    ("Life Events", "90%", "人生重要时刻没文字记录"),
    ("Search Photos", "12min", "平均找旧照片的时间"),
]
for i, (label, num, note) in enumerate(stats):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(3.2 + row * 2.2)
    tb(slide, x, y, Inches(3.5), Inches(0.2), label, Pt(8), INK_TERTIARY, font_name='SF Mono')
    tb(slide, x, y + Inches(0.25), Inches(3.5), Inches(0.6), num, Pt(42), INK, bold=True)
    tb(slide, x, y + Inches(0.95), Inches(3.5), Inches(0.3), note, Pt(12), INK_LIGHT)

# ============================================================
# SLIDE 03 · Quote + Concept (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "Act I · 记录方式的演变", True)
add_footer(slide, "历史与未来", "03 / 15", True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "从古至今", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(1.0),
   "编年史者，皇帝的 AI", Pt(48), PAPER, bold=True)
tb(slide, Inches(0.8), Inches(3.0), Inches(7), Inches(0.8),
   "古代帝王有起居注官。资治通鉴耗时19年。如果当时有AI，他们需要的只是一个App。",
   Pt(16), RGBColor(0xBB,0xBB,0xBB))

# Callout box
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.4), Inches(7), Inches(1.6))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1C)
shape.line.color.rgb = PAPER
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "记录不是为了别人，是为了未来的自己。"
p.font.size = Pt(18)
p.font.color.rgb = PAPER
p.font.name = 'PingFang SC'
p2 = tf.add_paragraph()
p2.text = "— HAVE DAY 的产品哲学"
p2.font.size = Pt(11)
p2.font.color.rgb = RGBColor(0x99,0x99,0x99)
p2.font.name = 'SF Mono'

# ============================================================
# SLIDE 04 · Act Divider (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "第二幕", False)
add_footer(slide, "第二幕引子 · 产品定义", "04 / 15", False)
tb(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.3),
   "Act II", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(1.2),
   "HAVE DAY", Pt(96), INK, bold=True)
tb(slide, Inches(0.8), Inches(3.8), Inches(9), Inches(0.6),
   "一款 AI 驱动的电子日记。你只需要说话、拍照，剩下的交给 AI。",
   Pt(20), INK_LIGHT)

# ============================================================
# SLIDE 05 · Pipeline (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "Act II · 产品功能", False)
add_footer(slide, "产品架构", "05 / 15", False)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "四大核心能力", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
   "HAVE DAY 的产品架构", Pt(48), INK, bold=True)

steps = [
    ("01", "语音 + 文字", "Web Speech API 实时转写，AI 自动润色修正"),
    ("02", "照片 + 记忆", "每天3张照片，自动压缩存储，AI 描述"),
    ("03", "日历追溯", "任何一天都能回溯，记忆永不丢失"),
    ("04", "AI 小哈", "你的专属日记助手——总结、安慰、鼓励"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.8 + i * 1.25)
    tb(slide, Inches(0.8), y, Inches(0.6), Inches(0.4), num, Pt(32), PRIMARY, bold=True)
    tb(slide, Inches(1.6), y, Inches(3), Inches(0.4), title, Pt(18), INK, bold=True)
    tb(slide, Inches(1.6), y + Inches(0.4), Inches(9), Inches(0.4), desc, Pt(13), INK_LIGHT)
    if i < 3:
        add_line(slide, Inches(1.6), y + Inches(1.0), Inches(10), INK_TERTIARY)

# ============================================================
# SLIDE 06 · Voice + AI (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "Act II · AI 文字润色", True)
add_footer(slide, "AI 润色 · 原文 vs 润色版", "06 / 15", True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "说出来的日记", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(1.0),
   "你说，AI 帮你写", Pt(48), PAPER, bold=True)
tb(slide, Inches(0.8), Inches(3.0), Inches(7), Inches(0.6),
   "对着手机说一段话，AI 自动修正错别字、补全语句、润色表达",
   Pt(16), RGBColor(0xBB,0xBB,0xBB))

# Example text
tb(slide, Inches(0.8), Inches(4.0), Inches(7), Inches(0.4),
   "原文：", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(4.3), Inches(7), Inches(0.5),
   "\"今天和同事去吃了那家新开的川菜馆，辣得我满头大汗但是超级爽\"",
   Pt(15), RGBColor(0xCC,0xCC,0xCC))
tb(slide, Inches(0.8), Inches(5.0), Inches(7), Inches(0.4),
   "AI 润色后：", Pt(11), PRIMARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(5.3), Inches(7), Inches(0.6),
   "\"今天和同事们探了一家新开川菜馆，辣得满头大汗却格外痛快。\"",
   Pt(15), PAPER)

# ============================================================
# SLIDE 07 · Image Grid (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "Act II · 功能一览", False)
add_footer(slide, "产品截图 · HAVE DAY v1.0", "07 / 15", False)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "From Text to Memory", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
   "三大核心交互", Pt(48), INK, bold=True)

features = [
    "📝 日记编辑器", "🎤 语音录制", "🐾 AI 小哈回复",
    "📷 照片上传", "📅 月历追溯", "📊 年度传记",
]
for i, feat in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(2.8 + row * 2.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.5), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE5, 0xDE)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = feat
    p.font.size = Pt(20)
    p.font.color.rgb = INK
    p.font.name = 'PingFang SC'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(20)

# ============================================================
# SLIDE 08 · Act Divider (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "第三幕", True)
add_footer(slide, "第三幕引子", "08 / 15", True)
tb(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.3),
   "Act III", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(1.2),
   "AI 不只是工具", Pt(84), PAPER, bold=True)
tb(slide, Inches(0.8), Inches(3.8), Inches(9), Inches(0.6),
   "它是你的倾听者、你的传记作家、你的生活搭子。",
   Pt(20), RGBColor(0xBB,0xBB,0xBB))

# ============================================================
# SLIDE 09 · Annual Report (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "Act III · 年度传记", False)
add_footer(slide, "年度传记 · 一键生成", "09 / 15", False)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "你的专属编年史", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(0.8),
   "一键生成年度传记", Pt(48), INK, bold=True)
tb(slide, Inches(0.8), Inches(2.8), Inches(7), Inches(0.6),
   "AI 读取全年日记，分析心情变化、提取高光时刻，用温暖的笔触写成你的年度故事。",
   Pt(16), INK_LIGHT)

# Callout
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(11), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE5, 0xDE)
shape.line.color.rgb = PRIMARY
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"亲爱的，这是你的 2026 年。这一年你记录了 217 天。春天学会了烘焙，夏天跑了第一次半马，秋天开始学吉他..."'
p.font.size = Pt(16)
p.font.color.rgb = INK
p.font.name = 'PingFang SC'
p2 = tf.add_paragraph()
p2.text = "— HAVE DAY 年度传记节选"
p2.font.size = Pt(10)
p2.font.color.rgb = INK_TERTIARY
p2.font.name = 'SF Mono'

# ============================================================
# SLIDE 10 · Alarm Reminder (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "Act III · 全屏提醒", True)
add_footer(slide, "提醒功能 · Apple 闹钟风格", "10 / 15", True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "像苹果闹钟一样", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(1.0),
   "重要的事，必须被看见", Pt(48), PAPER, bold=True)
tb(slide, Inches(0.8), Inches(3.0), Inches(7), Inches(0.6),
   "全屏黑色毛玻璃覆盖 + 长按3秒环形按钮确认。不是普通通知，是必须被确认的重要提醒。",
   Pt(16), RGBColor(0xBB,0xBB,0xBB))

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.4), Inches(7), Inches(1.6))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1C)
shape.line.color.rgb = PAPER
shape.line.width = Pt(1.5)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "提醒的本质不是'告知'，是'确认对方已经知道'。长按3秒就是那道确认的门槛。"
p.font.size = Pt(18)
p.font.color.rgb = PAPER
p.font.name = 'PingFang SC'
p2 = tf.add_paragraph()
p2.text = "— HAVE DAY 设计团队"
p2.font.size = Pt(11)
p2.font.color.rgb = RGBColor(0x99,0x99,0x99)
p2.font.name = 'SF Mono'

# ============================================================
# SLIDE 11 · Act Divider (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "第四幕", False)
add_footer(slide, "第四幕引子 · 技术架构", "11 / 15", False)
tb(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.3),
   "Act IV", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(1.2),
   "技术亮点", Pt(84), INK, bold=True)
tb(slide, Inches(0.8), Inches(3.8), Inches(9), Inches(0.6),
   "PWA 可安装、离线可用、本地存储、纯前端架构——打开浏览器就能用。",
   Pt(20), INK_LIGHT)

# ============================================================
# SLIDE 12 · Tech Stack (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "Act IV · 技术栈", True)
add_footer(slide, "技术选型 · 全部开源", "12 / 15", True)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "Under the Hood", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
   "技术栈", Pt(48), PAPER, bold=True)

tech_items = [
    ("前端框架", "React 18 + TypeScript + Vite 5"),
    ("样式 & 动效", "Tailwind CSS + Framer Motion（Apple 风格）"),
    ("本地存储", "Dexie.js (IndexedDB)——照片/日记/提醒持久化"),
    ("语音识别", "Web Speech API——浏览器原生，中文普通话"),
    ("AI 服务", "DeepSeek API——文本润色 + 日记总结 + 年报生成"),
    ("PWA 支持", "vite-plugin-pwa——可安装、离线浏览已有日记"),
    ("部署", "纯静态文件——Vercel / GitHub Pages / 阿里云 OSS"),
    ("数据隐私", "所有数据存储在用户手机浏览器本地——不上传服务器"),
]
for i, (label, value) in enumerate(tech_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(2.8 + row * 1.1)
    tb(slide, x, y, Inches(5.3), Inches(1.0),
       f"{label}：{value}", Pt(13), RGBColor(0xCC,0xCC,0xCC))

# ============================================================
# SLIDE 13 · Design (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "Act IV · 设计哲学", False)
add_footer(slide, "设计哲学 · Apple Minimalism", "13 / 15", False)
tb(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.3),
   "就像 iOS 原生应用", Pt(11), INK_TERTIARY, font_name='SF Mono')
tb(slide, Inches(0.8), Inches(1.7), Inches(7), Inches(0.8),
   "Apple 生态级的设计体验", Pt(48), INK, bold=True)
tb(slide, Inches(0.8), Inches(2.8), Inches(7), Inches(0.6),
   "毛玻璃效果(backdrop-filter)、Framer Motion spring动画、PingFang SC中文排版、safe-area刘海屏适配。",
   Pt(16), INK_LIGHT)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.0), Inches(11), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE5, 0xDE)
shape.line.color.rgb = PRIMARY
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"少即是多。每个页面只有一个核心动作，大量留白，让内容自己说话。"'
p.font.size = Pt(18)
p.font.color.rgb = INK
p.font.name = 'PingFang SC'
p2 = tf.add_paragraph()
p2.text = "— HAVE DAY 设计原则"
p2.font.size = Pt(10)
p2.font.color.rgb = INK_TERTIARY
p2.font.name = 'SF Mono'

# ============================================================
# SLIDE 14 · Big Quote (Dark)
# ============================================================
slide = add_dark_slide()
add_kicker(slide, "Act V · 愿景", True)
add_footer(slide, "金句 · 产品价值", "14 / 15", True)
tb(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.3),
   "Have A Nice Day", Pt(11), RGBColor(0x88,0x88,0x88), font_name='SF Mono', align=PP_ALIGN.CENTER)
tb(slide, Inches(1.0), Inches(2.5), Inches(11), Inches(2.0),
   '"当有一天，你忽然怀念和对象见第一面的那天——恰好你在 HAVE DAY 记录过。再次打开那天的日记，你会觉得这个 App 有价值。"',
   Pt(26), PAPER, align=PP_ALIGN.CENTER)
tb(slide, Inches(3), Inches(5.2), Inches(7), Inches(0.4),
   "— HAVE DAY 的初心 —", Pt(11), RGBColor(0x88,0x88,0x88),
   font_name='SF Mono', align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15 · Closing (Light)
# ============================================================
slide = add_light_slide()
add_kicker(slide, "HAVE DAY · Thank You", False)
add_footer(slide, "HAVE DAY · AI 电子日记", "— 2026 —", False)
tb(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.3),
   "谢谢", Pt(11), INK_TERTIARY, font_name='SF Mono', align=PP_ALIGN.CENTER)
tb(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(1.2),
   "HAVE DAY", Pt(96), INK, bold=True, align=PP_ALIGN.CENTER)
tb(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.8),
   "记录每一天，美好每一天。\n期待你成为 HAVE DAY 的第一批用户。",
   Pt(20), INK_LIGHT, align=PP_ALIGN.CENTER)
tb(slide, Inches(3), Inches(5.2), Inches(7), Inches(0.4),
   "have-day.vercel.app  ·  coming soon", Pt(11), INK_TERTIARY,
   font_name='SF Mono', align=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = '/Users/jk/Documents/Claud/have-day/ppt/HAVE_DAY.pptx'
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: {os.path.getsize(output_path) / 1024:.0f} KB")
